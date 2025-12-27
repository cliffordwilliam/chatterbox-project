import tomllib
import tempfile
from pathlib import Path
from chatterbox.tts_turbo import ChatterboxTurboTTS
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
import torchaudio as ta
from pptx.dml.color import RGBColor
from pygments import lex
from pygments.lexers import get_lexer_by_name
from pygments.styles import get_style_by_name
from pygments.token import Token
import tomli_w
import shutil


WIN_WORKDIR = Path("/mnt/c/Users/cclif/ppt-patcher")
WIN_INPUT_DIR = WIN_WORKDIR / "input"

PATCH_SCRIPT_SRC = Path(__file__).parent / "patch_pptx.py"
PATCH_SCRIPT_DST = WIN_WORKDIR / "patch_pptx.py"

PATCH_SCRIPT_SRC = Path(__file__).parent / "patch_pptx.py"
REQS_SRC = Path(__file__).parent / "requirements.txt"
CMD_SRC = Path(__file__).parent / "setup_and_patch.cmd"

PATCH_SCRIPT_DST = WIN_WORKDIR / "patch_pptx.py"
REQS_DST = WIN_WORKDIR / "requirements.txt"
CMD_DST = WIN_WORKDIR / "setup_and_patch.cmd"


def copy_windows_assets():
    for src, dst in [
        (PATCH_SCRIPT_SRC, PATCH_SCRIPT_DST),
        (REQS_SRC, REQS_DST),
        (CMD_SRC, CMD_DST),
    ]:
        if not src.exists():
            print(f"⚠ Missing file: {src}")
            continue

        shutil.copy2(src, dst)
        print(f"✓ Copied {src.name} → {dst}")


def ensure_windows_workspace():
    WIN_INPUT_DIR.mkdir(parents=True, exist_ok=True)


def hex_to_rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.lstrip("#")
    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16),
    )


def build_color_map(style_name="dracula"):
    style = get_style_by_name(style_name)
    color_map = {}

    for token, opts in style:
        if opts["color"]:
            color_map[token] = hex_to_rgb(opts["color"])

    return color_map


def add_code_with_highlighting(
    text_frame,
    code: str,
    lang: str,
    highlight_lines=None,
    font_size=10,
    font_name="Consolas",
):
    highlight_lines = set(highlight_lines or [])
    lexer = get_lexer_by_name(lang, stripall=False)
    color_map = build_color_map("dracula")

    text_frame.clear()

    for i, line in enumerate(code.splitlines(), start=1):
        p = text_frame.add_paragraph()
        p.line_spacing = 1.0

        is_highlight = i in highlight_lines

        for token_type, token_value in lex(line + "\n", lexer):
            run = p.add_run()
            run.text = token_value
            run.font.name = font_name
            run.font.size = Pt(font_size)

            t = token_type
            while t not in color_map and t is not Token:
                t = t.parent

            run.font.color.rgb = color_map.get(t, RGBColor(255, 255, 255))

            if highlight_lines and not is_highlight:
                run.font.color.brightness = -0.6


def generate_presentation(toml_path: str):
    ensure_windows_workspace()
    copy_windows_assets()

    with open(toml_path, "rb") as f:
        data = tomllib.load(f)

    slides_data = data["slides"]

    model = ChatterboxTurboTTS.from_pretrained(device="cuda")

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    for idx, slide_data in enumerate(slides_data):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 0, 0)

        left = Inches(0)
        top = Inches(0)
        width = prs.slide_width
        height = prs.slide_height

        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame

        text_frame.margin_left = Inches(0.25)
        text_frame.margin_right = Inches(0.25)
        text_frame.margin_top = Inches(0.25)
        text_frame.margin_bottom = Inches(0.25)
        text_frame.word_wrap = True

        add_code_with_highlighting(
            text_frame=text_frame,
            code=slide_data["text"],
            lang=slide_data.get("lang", "text"),
            highlight_lines=slide_data.get("highlight"),
        )

        wav = model.generate(slide_data["voice"])
        duration = wav.shape[1] / model.sr
        slide_data["duration"] = round(duration, 2)

        with tempfile.NamedTemporaryFile(suffix=".wav", delete=False) as temp_audio:
            temp_audio_path = temp_audio.name
            ta.save(temp_audio_path, wav, model.sr)

        slide.shapes.add_movie(
            temp_audio_path,
            left=Inches(0),
            top=Inches(0),
            width=Inches(1),
            height=Inches(1),
            mime_type="audio/wav",
        )

        Path(temp_audio_path).unlink()

        print(
            f"Processed slide {idx + 1}/{len(slides_data)} "
            f"(duration={slide_data['duration']}s)"
        )

    pptx_out = WIN_INPUT_DIR / "output.pptx"
    toml_out = WIN_INPUT_DIR / "presentation.toml"

    prs.save(pptx_out)

    with open(toml_out, "wb") as f:
        tomli_w.dump(data, f)

    print(f"\nPresentation saved to: {pptx_out}")
    print(f"PPTX written to: {pptx_out}")
    print(f"TOML written to: {toml_out}")


if __name__ == "__main__":
    generate_presentation("presentation.toml")
